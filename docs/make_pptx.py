from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0F, 0x1C, 0x3A)   # fondo oscuro
BLUE    = RGBColor(0x3B, 0x82, 0xF6)   # acento
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xCB, 0xD5, 0xE1)   # texto secundario
YELLOW  = RGBColor(0xFB, 0xBF, 0x24)   # highlight

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completamente en blanco

# ── Helpers ──────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    # fondo navy
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    return s

def box(s, x, y, w, h):
    return s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def line(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_after=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def accent_bar(s, y_inches=0.55):
    """Línea horizontal de acento."""
    line_shape = s.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.6), Inches(y_inches), Inches(1.2), Inches(0.06)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = BLUE
    line_shape.line.fill.background()

def badge(s, x, y, w, h, bg_color=BLUE):
    rect = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.fill.background()
    return rect

def pill_text(s, x, y, w, h, text, size=14, bg=BLUE, fg=WHITE):
    r = badge(s, x, y, w, h, bg)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Calibri"

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TÍTULO
# ─────────────────────────────────────────────────────────────────────────────
s1 = slide()

# Bloque de color izquierdo
left = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.8), H)
left.fill.solid()
left.fill.fore_color.rgb = BLUE
left.line.fill.background()

tb = box(s1, 0.4, 2.0, 4.0, 4.0)
tf = tb.text_frame
tf.word_wrap = True
line(tf, "Generador de", 32, bold=False, color=WHITE)
line(tf, "Oferta Académica", 40, bold=True,  color=WHITE, space_after=18)
line(tf, "Introducción al Desarrollo de Software", 13, color=RGBColor(0xE2,0xE8,0xF0))
line(tf, "Asistido por IA", 13, color=RGBColor(0xE2,0xE8,0xF0), space_after=30)
line(tf, "Agustín Maglione", 14, bold=True, color=WHITE)
line(tf, "CEIA — UBA", 12, color=RGBColor(0xE2,0xE8,0xF0))

# Lado derecho — keywords flotantes
kw_items = [
    (5.4, 1.2, "FastAPI"),
    (9.0, 1.0, "React 18"),
    (6.8, 2.6, "CP-SAT Optimizer"),
    (5.2, 3.9, "Multi-tenant"),
    (8.5, 3.4, "SQLAlchemy"),
    (6.5, 5.0, "OR-Tools"),
    (9.3, 5.2, "Tailwind"),
]
for x, y, label in kw_items:
    pill_text(s1, x, y, len(label)*0.13 + 0.5, 0.42, label, size=13,
              bg=RGBColor(0x1E,0x3A,0x6E), fg=LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — EL PROBLEMA
# ─────────────────────────────────────────────────────────────────────────────
s2 = slide()
accent_bar(s2)

tb = box(s2, 0.6, 0.7, 10, 0.9)
tf = tb.text_frame
line(tf, "El Problema", 30, bold=True, color=WHITE)

items = [
    ("📋", "Proceso manual y complejo cada semestre"),
    ("👨‍🏫", "Hay que cruzar docentes, franjas, aulas y correlatividades"),
    ("📊", "La demanda real depende del historial académico de cada alumno"),
    ("⚠️",  "Errores generan solapamientos, aulas sin uso o materias sin cubrir"),
]

for i, (icon, text) in enumerate(items):
    y = 1.7 + i * 1.2
    badge(s2, 0.6, y, 0.7, 0.7, RGBColor(0x1E,0x3A,0x6E))
    tb_icon = box(s2, 0.62, y + 0.04, 0.65, 0.65)
    tf_i = tb_icon.text_frame
    p = tf_i.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = icon; r.font.size = Pt(22)

    tb_t = box(s2, 1.55, y + 0.1, 10.5, 0.6)
    tf_t = tb_t.text_frame
    line(tf_t, text, 20, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — LA SOLUCIÓN (flujo)
# ─────────────────────────────────────────────────────────────────────────────
s3 = slide()
accent_bar(s3)

tb = box(s3, 0.6, 0.7, 10, 0.9)
tf = tb.text_frame
line(tf, "La Solución", 30, bold=True, color=WHITE)

steps = [
    ("Datos\nacadémicos", "Carreras, materias,\ndocentes, alumnos"),
    ("Demanda\nautomática", "Alumnos elegibles\npor correlatividad"),
    ("Optimizador\nCP-SAT", "Asigna docente\ny franja horaria"),
    ("Oferta\neditable", "Calendario con\nedición manual"),
    ("Insights", "Alertas y\nestadísticas"),
]

box_w = 2.1
gap   = 0.25
start_x = 0.35
y_box = 1.8

for i, (title, sub) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    rect = badge(s3, x, y_box, box_w, 2.2, RGBColor(0x1E,0x3A,0x6E))
    tf_r = rect.text_frame
    tf_r.word_wrap = True
    p1 = tf_r.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = title
    r1.font.size = Pt(17); r1.font.bold = True
    r1.font.color.rgb = BLUE; r1.font.name = "Calibri"

    p2 = tf_r.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(12); r2.font.color.rgb = LIGHT
    r2.font.name = "Calibri"

    # Flecha
    if i < len(steps) - 1:
        ax = x + box_w + 0.02
        arr = box(s3, ax, y_box + 0.85, gap + 0.1, 0.5)
        tf_a = arr.text_frame
        p_a = tf_a.paragraphs[0]
        p_a.alignment = PP_ALIGN.CENTER
        r_a = p_a.add_run(); r_a.text = "→"
        r_a.font.size = Pt(22); r_a.font.color.rgb = BLUE
        r_a.font.name = "Calibri"

# Tag multi-tenant
pill_text(s3, 0.6, 5.5, 2.8, 0.5, "✦  Multi-tenant", 14,
          bg=RGBColor(0x0A,0x3D,0x62), fg=YELLOW)
pill_text(s3, 3.8, 5.5, 3.4, 0.5, "✦  Background task async", 14,
          bg=RGBColor(0x0A,0x3D,0x62), fg=YELLOW)
pill_text(s3, 7.6, 5.5, 2.6, 0.5, "✦  59 tests pytest", 14,
          bg=RGBColor(0x0A,0x3D,0x62), fg=YELLOW)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — TECNOLOGÍAS
# ─────────────────────────────────────────────────────────────────────────────
s4 = slide()
accent_bar(s4)

tb = box(s4, 0.6, 0.7, 10, 0.9)
tf = tb.text_frame
line(tf, "Tecnologías", 30, bold=True, color=WHITE)

backend = [
    ("Python 3.11 + FastAPI", "API REST"),
    ("SQLAlchemy 2.0 + Alembic", "ORM + migraciones"),
    ("PostgreSQL", "Base de datos"),
    ("Google OR-Tools — CP-SAT", "Optimizador"),
    ("pytest", "59 tests automatizados"),
]
frontend = [
    ("React 18 + Vite", "Framework SPA"),
    ("Tailwind CSS + shadcn/ui", "Estilos + componentes"),
    ("axios + JWT", "HTTP client autenticado"),
    ("lucide-react", "Íconos"),
    ("sonner", "Notificaciones toast"),
]

col_labels = ["Backend", "Frontend"]
cols = [backend, frontend]
col_x = [0.55, 6.9]
col_w = 5.8

for ci, (col_x_val, col_items) in enumerate(zip(col_x, cols)):
    # Header de columna
    pill_text(s4, col_x_val, 1.55, 2.2, 0.48, col_labels[ci], 15,
              bg=BLUE, fg=WHITE)
    for row, (tech, desc) in enumerate(col_items):
        y = 2.25 + row * 0.88
        badge(s4, col_x_val, y, col_w, 0.72, RGBColor(0x16,0x28,0x4A))
        tb_r = box(s4, col_x_val + 0.18, y + 0.08, col_w - 0.3, 0.6)
        tf_r = tb_r.text_frame
        p_t = tf_r.paragraphs[0]
        rt = p_t.add_run(); rt.text = tech
        rt.font.size = Pt(15); rt.font.bold = True
        rt.font.color.rgb = WHITE; rt.font.name = "Calibri"
        p_d = tf_r.add_paragraph()
        rd = p_d.add_run(); rd.text = desc
        rd.font.size = Pt(11); rd.font.color.rgb = LIGHT
        rd.font.name = "Calibri"

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — MOTOR: REGLAS
# ─────────────────────────────────────────────────────────────────────────────
s5 = slide()
accent_bar(s5)

tb = box(s5, 0.6, 0.7, 12, 0.9)
tf = tb.text_frame
line(tf, "Motor de Optimización — CP-SAT", 30, bold=True, color=WHITE)

# Variables
tb_v = box(s5, 0.6, 1.55, 12, 0.55)
tf_v = tb_v.text_frame
line(tf_v, "Variable de decisión:  x[curso, docente, franja]  ∈  {0, 1}", 15,
     color=YELLOW)

hard = [
    "Cada curso → exactamente 1 (docente, franja)",
    "Un docente no puede dar 2 cursos en la misma franja",
    "Horas semanales por docente ≤ límite configurado",
    "Cursos simultáneos ≤ aulas disponibles",
    "Turnos permitidos por materia (allowed_turnos)",
]
soft = [
    "Minimizar solapamientos de materias del mismo año/carrera",
    "Balancear carga horaria entre franjas",
]

# Columna izquierda — duras
pill_text(s5, 0.55, 2.25, 2.8, 0.42, "RESTRICCIONES DURAS", 12,
          bg=RGBColor(0xDC,0x26,0x26), fg=WHITE)
for i, h in enumerate(hard):
    tb_h = box(s5, 0.7, 2.85 + i * 0.75, 5.7, 0.62)
    tf_h = tb_h.text_frame
    line(tf_h, f"  {h}", 14, color=WHITE)

# Columna derecha — blandas
pill_text(s5, 7.1, 2.25, 2.6, 0.42, "OBJETIVO (blandas)", 12,
          bg=RGBColor(0xD9,0x77,0x06), fg=WHITE)
for i, so in enumerate(soft):
    tb_s = box(s5, 7.25, 2.85 + i * 0.75, 5.5, 0.62)
    tf_s = tb_s.text_frame
    line(tf_s, f"  {so}", 14, color=WHITE)

# Timeout note
tb_t2 = box(s5, 0.6, 6.6, 12, 0.5)
tf_t2 = tb_t2.text_frame
line(tf_t2, "Si el solver agota el tiempo → devuelve la mejor solución encontrada (status: feasible, no optimal)",
     12, color=LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — CÓMO AYUDÓ CLAUDE
# ─────────────────────────────────────────────────────────────────────────────
s6 = slide()
accent_bar(s6)

tb = box(s6, 0.6, 0.7, 12, 0.9)
tf = tb.text_frame
line(tf, "Cómo ayudó Claude Code", 30, bold=True, color=WHITE)

# Ciclo de skills como flow
skills = [
    ("/brainstorming",              "Diseño colaborativo\nqué construir y cómo"),
    ("/writing-plans",              "Plan TDD paso a paso\ncon código exacto"),
    ("/subagent-driven-dev",        "Implementer + Reviewer\nfrescos por tarea"),
    ("/requesting-code-review",     "Review final\nde rama completa"),
    ("/finishing-a-dev-branch",     "Tests → push/PR\ncontrolado"),
]

bw = 2.35
by = 1.7
bh = 2.0
for i, (skill, desc) in enumerate(skills):
    bx = 0.3 + i * (bw + 0.1)
    rect = badge(s6, bx, by, bw, bh, RGBColor(0x1E,0x3A,0x6E))
    tf_r = rect.text_frame
    tf_r.word_wrap = True

    p1 = tf_r.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(10)
    r1 = p1.add_run(); r1.text = skill
    r1.font.size = Pt(12); r1.font.bold = True
    r1.font.color.rgb = BLUE; r1.font.name = "Calibri"

    p2 = tf_r.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(12); r2.font.color.rgb = LIGHT
    r2.font.name = "Calibri"

    if i < len(skills) - 1:
        ax = bx + bw
        arr = box(s6, ax, by + 0.7, 0.15, 0.5)
        tf_a = arr.text_frame
        p_a = tf_a.paragraphs[0]
        p_a.alignment = PP_ALIGN.CENTER
        r_a = p_a.add_run(); r_a.text = "→"
        r_a.font.size = Pt(18); r_a.font.color.rgb = BLUE
        r_a.font.name = "Calibri"

# Rol del humano
tb_h = box(s6, 0.6, 4.05, 12, 0.55)
tf_h = tb_h.text_frame
line(tf_h, "Rol del humano: definir QUÉ construir, elegir entre opciones, aprobar cada etapa.",
     15, color=YELLOW)

# Ejemplos de bugs encontrados por review
bug_tb = box(s6, 0.6, 4.75, 12, 2.0)
tf_b = bug_tb.text_frame
line(tf_b, "Bugs reales detectados por los subagentes de review:", 13, bold=True, color=LIGHT)
bugs = [
    'Modelo CP-SAT vacío devolvía "optimal" cuando todas las materias tenían turno restringido',
    'Turnos se ordenaban alfabéticamente: "Noche" antes que "Tarde" en estadísticas',
    "TurnosBadge no distinguía null (todos los turnos) de [] (ningún turno)",
]
for b in bugs:
    line(tf_b, f"  ·  {b}", 13, color=WHITE, space_after=2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — DECISIONES DE DISEÑO
# ─────────────────────────────────────────────────────────────────────────────
s7 = slide()
accent_bar(s7)

tb = box(s7, 0.6, 0.7, 12, 0.9)
tf = tb.text_frame
line(tf, "Decisiones de Diseño", 30, bold=True, color=WHITE)

decisions = [
    ("CP-SAT vs heurísticas",
     "Optimalidad garantizada o falla explícita con razón concreta"),
    ("Background task async",
     "El solver puede tardar minutos — la API no bloquea"),
    ("Multi-tenant desde el inicio",
     "Agregar multi-tenancy retroactivamente es costoso y riesgoso"),
    ("Insights en el backend",
     "El backend tiene nombres resueltos; el frontend solo renderiza"),
    ("Partial infeasibility",
     'Reportar "Física II sin turno válido" es mas util que "la generacion fallo"'),
    ("localStorage para sesión",
     "UX básica: recargar la página no debería forzar un nuevo login"),
]

cols_d = 2
rows_d = 3
card_w = 5.8
card_h = 1.55
gap_x  = 0.45
gap_y  = 0.25
start_x = 0.55
start_y = 1.65

for i, (title, desc) in enumerate(decisions):
    col = i % cols_d
    row = i // cols_d
    cx = start_x + col * (card_w + gap_x)
    cy = start_y + row * (card_h + gap_y)

    rect = badge(s7, cx, cy, card_w, card_h, RGBColor(0x16,0x28,0x4A))
    # Borde izquierdo de acento
    accent = badge(s7, cx, cy, 0.08, card_h, BLUE)

    tb_c = box(s7, cx + 0.2, cy + 0.15, card_w - 0.3, card_h - 0.2)
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    line(tf_c, title, 15, bold=True, color=WHITE, space_after=4)
    line(tf_c, desc,  13, color=LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — PRÓXIMOS PASOS
# ─────────────────────────────────────────────────────────────────────────────
s8 = slide()
accent_bar(s8)

tb = box(s8, 0.6, 0.7, 12, 0.9)
tf = tb.text_frame
line(tf, "Próximos Pasos", 30, bold=True, color=WHITE)

nexts = [
    ("🔔", "Notificaciones en tiempo real",
           "WebSocket para avisar cuando el job de generación termina"),
    ("👥", "Gestión de alumnos desde la UI",
           "Carga de historial académico sin tocar la base directamente"),
    ("📊", "Comparación entre versiones de oferta",
           "Ver diferencias antes de aprobar una nueva generación"),
    ("📄", "Export a formatos institucionales",
           "PDF y Excel para integración con sistemas académicos existentes"),
    ("🎯", "Restricciones de carga por carrera",
           "Hoy el límite de horas es global; permitir configurarlo por carrera"),
]

for i, (icon, title, desc) in enumerate(nexts):
    y = 1.6 + i * 1.0
    # Número
    pill_text(s8, 0.55, y + 0.05, 0.55, 0.55, icon, 18,
              bg=RGBColor(0x1E,0x3A,0x6E), fg=WHITE)
    tb_n = box(s8, 1.3, y + 0.0, 11.2, 0.85)
    tf_n = tb_n.text_frame
    tf_n.word_wrap = True
    line(tf_n, title, 17, bold=True, color=WHITE, space_after=2)
    line(tf_n, desc,  13, color=LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — DEMO / CIERRE
# ─────────────────────────────────────────────────────────────────────────────
s9 = slide()

left2 = s9.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.8), H)
left2.fill.solid()
left2.fill.fore_color.rgb = BLUE
left2.line.fill.background()

tb_d = box(s9, 0.5, 2.5, 3.8, 3.0)
tf_d = tb_d.text_frame
line(tf_d, "Demo", 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
line(tf_d, "", 12)
line(tf_d, "github.com/amaglione/", 13, color=RGBColor(0xE2,0xE8,0xF0),
     align=PP_ALIGN.CENTER)
line(tf_d, "academic-offer-generator", 13, color=RGBColor(0xE2,0xE8,0xF0),
     align=PP_ALIGN.CENTER)

summary_items = [
    "Sistema web end-to-end completo",
    "Motor CP-SAT con 5 restricciones duras",
    "Calendario editable con drag & drop",
    "Insights post-generación con alertas",
    "59 tests automatizados",
    "Desarrollado 100% con Claude Code",
]

tb_s = box(s9, 5.3, 1.5, 7.5, 5.5)
tf_s = tb_s.text_frame
tf_s.word_wrap = True
for item in summary_items:
    line(tf_s, f"  ✓  {item}", 18, color=WHITE, space_after=8)

OUT = "/Users/amaglione/Documents/Trabajos/ceia-uba/academic-offer-generator/docs/presentacion.pptx"
prs.save(OUT)
print(f"Guardado en {OUT}")
